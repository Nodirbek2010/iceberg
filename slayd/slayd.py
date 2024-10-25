from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1 - Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Удивительный талантливый ребёнок: Танишк Абрахам"
subtitle.text = "История о юном вундеркинде, изменяющем науку"

# Slide 2 - Образование и достижения
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
content_2 = slide_2.shapes.placeholders[1].text_frame

title_2.text = "Образование и успехи в юном возрасте"
content_2.text = "Танишк Абрахам начал демонстрировать свои способности с ранних лет.\n"
content_2.add_paragraph("• В 10 лет он закончил среднюю школу.")
content_2.add_paragraph("• В 11 лет поступил в колледж и стал самым молодым выпускником.")
content_2.add_paragraph("• В 14 лет поступил в Калифорнийский университет в Дэвисе.")
content_2.add_paragraph("• Танишк расширял свои знания через онлайн-курсы и платформы.")

# Slide 3 - Научные достижения
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
content_3 = slide_3.shapes.placeholders[1].text_frame

title_3.text = "Научные проекты и исследования"
content_3.text = "• Участвовал в исследовательских проектах NASA.\n"
content_3.add_paragraph("• Представил свою первую научную работу в 9 лет на тему глобального потепления.")
content_3.add_paragraph("• Исследовал технологии для переселения людей на Марс.")
content_3.add_paragraph("• Один из самых юных участников научных конференций.")

# Slide 4 - Личная жизнь и интересы
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
content_4 = slide_4.shapes.placeholders[1].text_frame

title_4.text = "Личная жизнь, хобби и общественная деятельность"
content_4.text = "• Танишк любит читать книги, заниматься спортом и проводить время с семьёй.\n"
content_4.add_paragraph("• Активно участвует в научных форумах и проектах для молодых учёных.")
content_4.add_paragraph("• Увлекается программированием и создаёт собственные научные приложения.")

# Slide 5 - Вдохновляющий герой
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
content_5 = slide_5.shapes.placeholders[1].text_frame

title_5.text = "Танишк как пример для молодёжи"
content_5.text = "• История Танишка вдохновляет многих молодых людей по всему миру.\n"
content_5.add_paragraph("• Его стремление к знаниям показывает, что возраст не помеха для достижения высот.")
content_5.add_paragraph("• Он стал примером для детей, мотивируя их стремиться к мечтам.")
content_5.add_paragraph("• Его поддержка и участие в STEM-проектах помогает другим детям развиваться.")

# Slide 6 - Заключение
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
content_6 = slide_6.shapes.placeholders[1].text_frame

title_6.text = "Заключение: Пути успеха и будущее Танишка"
content_6.text = "• Танишк продолжает развиваться и вносить вклад в науку и технологии.\n"
content_6.add_paragraph("• Его история показывает важность поддержки детей в их начинаниях.")
content_6.add_paragraph("• В будущем Танишк мечтает стать учёным и сделать открытия для всего человечества.")
content_6.add_paragraph("• Его успехи — пример для того, как дети могут изменить мир.")

# Save the presentation
pptx_file = "/mnt/data/Tanishk_Abraham_presentation.pptx"
prs.save(pptx_file)

pptx_file